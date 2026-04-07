from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image

def add_icon(slide, icon_path, left, top, width, label=None):
    pic = slide.shapes.add_picture(icon_path, Inches(left), Inches(top), width=Inches(width))
    if label:
        tx = slide.shapes.add_textbox(Inches(left), Inches(top + width + 0.05),
                                      Inches(width), Inches(0.3))
        tf = tx.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = 1  # center
    return pic

def add_box(slide, title, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)
    shape.line.color.rgb = RGBColor(120,120,120)
    shape.text = title
    shape.text_frame.paragraphs[0].font.size = Pt(12)
    shape.text_frame.paragraphs[0].font.bold = True
    return shape

# --- Presentation setup ---
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

# Region / VPC boxes
add_box(slide, "Thailand Region", 0.5, 0.8, 13.0, 7.0, (232,243,254))
add_box(slide, "POS Production VPC", 1.0, 1.2, 12.0, 6.0, (255,255,255))
add_box(slide, "ECS Cluster", 2.0, 1.5, 10.0, 5.3, (255,244,229))

# Subnets
add_box(slide, "Ingress subnet", 1.2, 1.6, 2.2, 1.7, (247,247,247))
add_box(slide, "App Tier subnet", 3.0, 1.6, 2.5, 4.7, (239,255,242))
add_box(slide, "DB Tier subnet", 5.7, 1.6, 2.5, 4.7, (255,241,243))
add_box(slide, "Private subnet EP", 8.4, 1.6, 2.5, 4.7, (237,248,255))

# --- Main icons ---
icon_dir = "aws_icons_2023"

add_icon(slide, f"{icon_dir}/API-Gateway.png", 1.9, 1.8, 0.6, "API Gateway")
add_icon(slide, f"{icon_dir}/Application-Load-Balancer.png", 2.6, 1.8, 0.6, "ALB / SG")
add_icon(slide, f"{icon_dir}/ECS.png", 3.4, 2.0, 0.7, "Container")
add_icon(slide, f"{icon_dir}/ElastiCache-Redis.png", 6.2, 2.3, 0.6, "Redis")
add_icon(slide, f"{icon_dir}/RDS-PostgreSQL.png", 6.2, 3.5, 0.6, "PostgreSQL")
add_icon(slide, f"{icon_dir}/S3.png", 10.0, 3.0, 0.7, "S3 Bucket")
add_icon(slide, f"{icon_dir}/WAF.png", 2.5, 7.1, 0.6, "WAF")
add_icon(slide, f"{icon_dir}/ECR.png", 3.2, 7.1, 0.6, "ECR")
add_icon(slide, f"{icon_dir}/Secrets-Manager.png", 3.9, 7.1, 0.6, "Secrets Manager")
add_icon(slide, f"{icon_dir}/Systems-Manager-Parameter-Store.png", 4.6, 7.1, 0.6, "Parameter Store")
add_icon(slide, f"{icon_dir}/CloudWatch.png", 5.3, 7.1, 0.6, "CloudWatch Logs")
add_icon(slide, f"{icon_dir}/Backup.png", 6.0, 7.1, 0.6, "AWS Backup")
add_icon(slide, f"{icon_dir}/Certificate-Manager.png", 6.7, 7.1, 0.6, "Certificate Manager")
add_icon(slide, f"{icon_dir}/Transit-Gateway.png", 0.6, 4.5, 0.7, "TGW")
add_icon(slide, f"{icon_dir}/Firewall-Manager.png", 0.6, 3.3, 0.7, "Palo FW")

# --- Output ---
prs.save("Cloud_Example_Diagram.pptx")
print("✅  Cloud_Example_Diagram.pptx created successfully.")
